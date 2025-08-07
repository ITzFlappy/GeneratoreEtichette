import importlib.metadata
import streamlit as st
import pandas as pd
from pptx import Presentation
from io import BytesIO
import tempfile
from copy import deepcopy
import re
from pathlib import Path
from datetime import datetime

# Cartella dove metti i template .pptx
TEMPLATE_DIR = Path("Templates")
available_templates = {f.stem: f for f in TEMPLATE_DIR.glob("*.pptx")}

# Rimuove il prefisso 'NTemplate' dalla visualizzazione
def clean_name(name):
    return name.removeprefix("NTemplate")

template_display_names = [clean_name(name) for name in available_templates]
display_to_filename = {clean_name(name): name for name in available_templates}

# Evita errore su importlib.metadata.version
importlib.metadata.version = lambda name: "1.48.0" if name == "streamlit" else importlib.metadata.version(name)

def replace_text_in_shapes(slide, data_dict):
    pattern = r"<(.*?)>"
    for shape in slide.shapes:
        if shape.has_text_frame:
            full_text = ""
            runs = []
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    runs.append(run)
                    full_text += run.text

            matches = re.findall(pattern, full_text)
            for key in matches:
                placeholder = f"<{key}>"
                replacement = str(data_dict.get(key.strip(), ""))
                full_text = full_text.replace(placeholder, replacement)

            # Sovrascrivi tutto il testo nel primo run
            if runs:
                runs[0].text = full_text
                for r in runs[1:]:
                    r.text = ""  # pulisci i run successivi
    return slide

def duplicate_slide(prs, source_slide):
    slide_layout = prs.slide_layouts[0]
    new_slide = prs.slides.add_slide(slide_layout)
    for shape in source_slide.shapes:
        try:
            if shape.shape_type == 13:
                image_stream = BytesIO(shape.image.blob)
                new_slide.shapes.add_picture(
                    image_stream, shape.left, shape.top, shape.width, shape.height
                )
            else:
                el = shape.element
                new_el = deepcopy(el)
                new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
        except Exception:
            pass
    return new_slide

st.set_page_config(page_title="Generatore Etichette", layout="centered")
st.title("Generatore Etichette")

search_term = st.text_input("Cerca un template", "")
filtered_display_names = [name for name in template_display_names if search_term.lower() in name.lower()]
selected_display_name = st.selectbox("Seleziona un template", filtered_display_names)

# Nome del file template effettivo
selected_template_stem = display_to_filename.get(selected_display_name)
excel_file = st.file_uploader("Carica il file Excel (.xlsx o .xls)", type=["xlsx", "xls"])

if selected_template_stem and excel_file:
    df = pd.read_excel(excel_file)
    st.write("Colonne Excel:", list(df.columns))
    st.success(f"{len(df)} righe caricate dal file Excel.")

    if st.button("Genera PowerPoint"):
        selected_path = available_templates[selected_template_stem]
        template_ppt = Presentation(str(selected_path))
        template_slide = template_ppt.slides[0]

        final_ppt = Presentation()
        final_ppt.slides._sldIdLst.clear()
        final_ppt.slide_width = template_ppt.slide_width
        final_ppt.slide_height = template_ppt.slide_height

        for _, row in df.iterrows():
            data_dict = {
                str(k).strip(): "" if pd.isna(v) else str(v).strip()
                for k, v in row.items()
            }
            new_slide = duplicate_slide(final_ppt, template_slide)
            replace_text_in_shapes(new_slide, data_dict)

        with BytesIO() as result:
            final_ppt.save(result)
            result.seek(0)
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"{selected_display_name}_{timestamp}.pptx"
            st.download_button(
                "Scarica PowerPoint Compilato",
                data=result,
                file_name=filename,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
